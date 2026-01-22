"""
template_engine.py - CV PPT 範本引擎

使用 CV 標準範本.pptx 作為範本，生成高階主管 CV 簡報。
"""

import io
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .field_formatter import FIELD_CONFIG, format_field_content


class CVTemplateEngine:
    """CV PPT 範本引擎"""

    # 範本檔案名稱（依優先順序）
    TEMPLATE_PATH = "CV 標準範本.pptx"
    TEMPLATE_PATH_ALT = "CV_標準範本.pptx"
    TEMPLATE_PATH_FALLBACK = "CV 範本.pptx"

    # 左側文字框欄位（依序）
    LEFT_FIELDS = ["專業背景", "學歷", "主要經歷"]

    # 右側文字框欄位（依序）
    RIGHT_FIELDS = ["現任", "個人特質", "現擔任獨董家數", "擔任獨董年資"]

    # Shape 位置常數（英吋）
    LEFT_TEXTBOX_POSITION = 2.74  # Shape 3
    RIGHT_TEXTBOX_POSITION = 7.46  # Shape 4

    def __init__(self):
        self.prs = None
        self.slide = None
        self.template_used = None

    def load_template(self) -> bool:
        """
        載入 CV 範本檔案

        Returns:
            是否成功載入
        """
        # 依優先順序嘗試載入範本
        template_paths = [
            self.TEMPLATE_PATH,
            self.TEMPLATE_PATH_ALT,
            self.TEMPLATE_PATH_FALLBACK
        ]

        self.template_used = None
        for path in template_paths:
            if Path(path).exists():
                self.template_used = path
                break

        if self.template_used is None:
            return False

        try:
            self.prs = Presentation(self.template_used)
            if self.prs.slides:
                self.slide = self.prs.slides[0]
            return True
        except Exception as e:
            return False

    def set_name(self, name: str, age: str = None):
        """
        設定姓名和年齡（Shape 1 - PLACEHOLDER）

        Args:
            name: 姓名字串
            age: 年齡字串（可選，如 "52歲"）
        """
        if not self.slide:
            return

        # 找到第一個 placeholder 或文字框（通常是姓名欄位）
        for shape in self.slide.shapes:
            if shape.has_text_frame:
                # 檢查是否是標題位置（通常在投影片上方）
                if shape.top < Inches(1.0):  # 位置在上方 1 英吋以內
                    self._set_name_with_age(shape, name, age)
                    return

        # 如果沒有找到，嘗試找 placeholder
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if shape.has_text_frame:
                    self._set_name_with_age(shape, name, age)
                    return

    def _set_name_with_age(self, shape, name: str, age: str = None):
        """設定姓名和年齡"""
        if not shape.has_text_frame:
            return

        tf = shape.text_frame

        # 清空所有段落
        for para in tf.paragraphs:
            para.clear()

        # 填入姓名
        p = tf.paragraphs[0]
        run_name = p.add_run()
        run_name.text = name
        run_name.font.name = "微軟正黑體"
        run_name.font.size = Pt(28)
        run_name.font.bold = True

        # 如果有年齡，加在後面
        if age and str(age).strip() and str(age).lower() not in ['nan', 'none', '']:
            run_age = p.add_run()
            # 確保年齡格式正確（加上括號）
            age_str = str(age).strip()
            if not age_str.startswith('（') and not age_str.startswith('('):
                age_str = f"（{age_str}）"
            run_age.text = f" {age_str}"
            run_age.font.name = "微軟正黑體"
            run_age.font.size = Pt(18)
            run_age.font.bold = False

    def _set_shape_text(self, shape, text: str):
        """設定形狀的文字內容（保留向後相容）"""
        if not shape.has_text_frame:
            return

        # 清空並設定新文字
        tf = shape.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = text
            run.font.name = "微軟正黑體"
            run.font.size = Pt(24)
            run.font.bold = True

    def set_age(self, age: str):
        """
        設定年齡（Shape 5 Group 內的子形狀）

        Args:
            age: 年齡字串，如 "52歲"。若為空/None/nan，則顯示「年齡：(待確認)」
        """
        if not self.slide:
            return

        # 找到 Shape 5 (GROUP) - 位置約 left=0.44", top=3.47"
        group_shape = None
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 確認是年齡的 Group（位置約 left=0.44", top=3.47"）
                left_inches = shape.left / 914400  # EMU to inches
                top_inches = shape.top / 914400
                if 0.3 < left_inches < 0.6 and 3.0 < top_inches < 4.0:
                    group_shape = shape
                    break

        if group_shape is None:
            print("  警告：找不到年齡 Group")
            return

        # 在 Group 內找到包含「年齡」的子形狀
        for sub_shape in group_shape.shapes:
            if sub_shape.has_text_frame:
                full_text = ''.join(run.text for para in sub_shape.text_frame.paragraphs for run in para.runs)
                if '年齡' in full_text:
                    # 清空並重新填入
                    tf = sub_shape.text_frame
                    for para in tf.paragraphs:
                        para.clear()

                    p = tf.paragraphs[0]
                    run = p.add_run()

                    # 處理年齡值
                    if age and str(age).strip() and str(age).lower() not in ['nan', 'none', '']:
                        age_str = str(age).strip()
                        # 確保格式正確
                        if not age_str.endswith('歲'):
                            age_str = f"{age_str}歲"
                        run.text = f"年齡：{age_str}"
                    else:
                        run.text = "年齡：(待確認)"

                    run.font.name = "微軟正黑體"
                    run.font.size = Pt(11)
                    run.font.bold = False

                    print(f"  → 設定年齡: {run.text}")
                    return

        print("  警告：在 Group 中找不到年齡欄位")

    # 照片預設位置和大小（英吋）
    PHOTO_DEFAULT_LEFT = 0.44
    PHOTO_DEFAULT_TOP = 0.85
    PHOTO_DEFAULT_WIDTH = 2.0
    PHOTO_DEFAULT_HEIGHT = 2.5

    def set_photo(self, image_source) -> bool:
        """
        設定照片（如果範本有 PICTURE 則替換，否則新增）

        Args:
            image_source: 圖片來源，可以是檔案路徑或 BytesIO

        Returns:
            是否成功設定
        """
        if not self.slide:
            return False

        if image_source is None:
            return False

        # 找到圖片形狀
        picture_shape = None
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shape = shape
                break

        if picture_shape is not None:
            # 範本有圖片 - 替換模式
            left = picture_shape.left
            top = picture_shape.top
            width = picture_shape.width
            height = picture_shape.height

            # 移除原始圖片
            sp = picture_shape._element
            sp.getparent().remove(sp)
        else:
            # 範本沒有圖片 - 新增模式（使用預設位置）
            left = Inches(self.PHOTO_DEFAULT_LEFT)
            top = Inches(self.PHOTO_DEFAULT_TOP)
            width = Inches(self.PHOTO_DEFAULT_WIDTH)
            height = Inches(self.PHOTO_DEFAULT_HEIGHT)
            print("  → 範本無圖片形狀，使用預設位置新增照片")

        # 處理圖片來源
        try:
            if isinstance(image_source, io.BytesIO):
                image_source.seek(0)
                self.slide.shapes.add_picture(image_source, left, top, width, height)
            elif isinstance(image_source, (str, Path)):
                if Path(image_source).exists():
                    self.slide.shapes.add_picture(str(image_source), left, top, width, height)
                else:
                    return False
            return True
        except Exception as e:
            print(f"  警告：新增照片失敗 - {e}")
            return False

    def fill_left_content(self, data: dict):
        """
        填充左側文字框（Shape 3）

        Args:
            data: {
                "專業背景": "...",
                "學歷": "1. 台大\\n2. 政大",
                "主要經歷": "1. NVIDIA\\n2. Microsoft"
            }
        """
        if not self.slide:
            return

        # 找到左側文字框（根據位置）
        shape = self._find_shape_by_position(self.LEFT_TEXTBOX_POSITION)
        if shape:
            self._fill_textbox(shape, self.LEFT_FIELDS, data)

    def fill_right_content(self, data: dict):
        """
        填充右側文字框（Shape 4）

        Args:
            data: {
                "現任": "1. 台大教授\\n2. 協會理事",
                "個人特質": "...",
                "現擔任獨董家數": 0,
                "擔任獨董年資": 0
            }
        """
        if not self.slide:
            return

        # 找到右側文字框（根據位置）
        shape = self._find_shape_by_position(self.RIGHT_TEXTBOX_POSITION)
        if shape:
            self._fill_textbox(shape, self.RIGHT_FIELDS, data)

    def _find_shape_by_position(self, target_left_inches: float, tolerance: float = 0.5):
        """
        根據 left 位置找到對應的 Shape

        Args:
            target_left_inches: 目標 left 位置（英吋）
            tolerance: 容許誤差（英吋）

        Returns:
            找到的 Shape 或 None
        """
        if not self.slide:
            return None

        target_left = Inches(target_left_inches)
        tolerance_emu = Inches(tolerance)

        for shape in self.slide.shapes:
            if shape.has_text_frame:
                # 檢查位置是否匹配
                if abs(shape.left - target_left) < tolerance_emu:
                    return shape

        return None

    def _calculate_content_metrics(self, fields: list, data: dict) -> dict:
        """
        計算內容指標，用於決定字體大小

        Returns:
            {
                "total_lines": 總行數,
                "total_chars": 總字元數,
                "is_dense": 是否內容密集
            }
        """
        total_lines = 0
        total_chars = 0

        for field_name in fields:
            config = FIELD_CONFIG.get(field_name, {})
            excel_col = config.get("excel_column")
            if excel_col and excel_col in data:
                raw_value = data.get(excel_col, "")
            else:
                raw_value = data.get(field_name, "")

            if raw_value is None or (isinstance(raw_value, float) and pd.isna(raw_value)):
                raw_value = ""

            content = str(raw_value).strip()
            if not content:
                content = config.get("empty_text", "(待補充)")

            # 計算行數
            lines = content.split("\n")
            total_lines += len([l for l in lines if l.strip()])
            total_chars += len(content)

        return {
            "total_lines": total_lines,
            "total_chars": total_chars,
            "is_dense": total_lines > 12 or total_chars > 400
        }

    def _get_adaptive_font_sizes(self, metrics: dict) -> dict:
        """
        根據內容指標返回自適應字體大小

        Returns:
            {
                "title": 標題字體大小,
                "content": 內容字體大小,
                "sub_item": 子項目字體大小
            }
        """
        total_lines = metrics["total_lines"]
        total_chars = metrics["total_chars"]

        # 根據內容密度調整字體
        if total_lines > 18 or total_chars > 600:
            # 非常密集 - 最小字體
            return {"title": Pt(9), "content": Pt(9), "sub_item": Pt(8)}
        elif total_lines > 15 or total_chars > 500:
            # 密集 - 小字體
            return {"title": Pt(9.5), "content": Pt(9.5), "sub_item": Pt(8.5)}
        elif total_lines > 12 or total_chars > 400:
            # 中等密集 - 較小字體
            return {"title": Pt(10), "content": Pt(10), "sub_item": Pt(9)}
        elif total_lines > 9 or total_chars > 300:
            # 輕微密集 - 稍小字體
            return {"title": Pt(11), "content": Pt(11), "sub_item": Pt(10)}
        else:
            # 正常 - 標準字體
            return {"title": Pt(12), "content": Pt(12), "sub_item": Pt(10.5)}

    def _fill_textbox(self, shape, fields: list, data: dict):
        """
        清空文字框並填入格式化內容（支援自動調整字體大小）

        格式規則：
        - 標題獨立一行（如「學歷：」後面換行）
        - 每筆資料獨立一行
        - 不同欄位之間空一行
        - 不要有多餘的空行

        Args:
            shape: pptx 形狀物件
            fields: 欄位名稱列表
            data: 資料字典
        """
        if not shape.has_text_frame:
            return

        tf = shape.text_frame

        # === 步驟 1：真正清空 text_frame ===
        # python-pptx 的 clear() 可能不完整，需要手動處理
        # 刪除第一個以外的所有段落
        while len(tf.paragraphs) > 1:
            p = tf.paragraphs[-1]
            p._element.getparent().remove(p._element)

        # 清空第一個段落
        first_para = tf.paragraphs[0]
        first_para.clear()

        # === 步驟 2：填入內容 ===
        is_first_field = True

        for field_name in fields:
            config = FIELD_CONFIG.get(field_name, {
                "label": field_name,
                "excel_column": None,
                "multiline": False,
                "empty_text": "(待補充)"
            })

            # 取得欄位值（優先使用 excel_column 映射）
            excel_col = config.get("excel_column")
            if excel_col and excel_col in data:
                raw_value = data.get(excel_col, "")
            else:
                raw_value = data.get(field_name, "")

            # 處理空值
            if raw_value is None or (isinstance(raw_value, float) and pd.isna(raw_value)) or str(raw_value).strip() == "":
                raw_value = config.get("empty_text", "(待補充)")

            content = str(raw_value).strip()

            # === 欄位之間空一行（除了第一個欄位）===
            if not is_first_field:
                p_spacer = tf.add_paragraph()
                run_spacer = p_spacer.add_run()
                run_spacer.text = ""

            # === 標題段落（獨立一行）===
            if is_first_field:
                p_title = first_para
                is_first_field = False
            else:
                p_title = tf.add_paragraph()

            # 加入粗體標題和冒號（12pt 粗體）
            run_title = p_title.add_run()
            run_title.text = f"{config.get('label', field_name)}："
            run_title.font.bold = True
            run_title.font.name = "微軟正黑體"
            run_title.font.size = Pt(12)

            # === 內容段落 ===
            if config.get("multiline", False) and "\n" in content:
                # 多行內容：每行一個段落
                lines = [line.strip() for line in content.split("\n") if line.strip()]
                for line in lines:
                    p_content = tf.add_paragraph()
                    run_content = p_content.add_run()
                    run_content.text = line
                    run_content.font.bold = False
                    run_content.font.name = "微軟正黑體"
                    # 所有內文（包含子項目）統一使用 11pt
                    run_content.font.size = Pt(11)
            else:
                # 單行內容：接在標題下一行（11pt）
                p_content = tf.add_paragraph()
                run_content = p_content.add_run()
                run_content.text = content
                run_content.font.bold = False
                run_content.font.name = "微軟正黑體"
                run_content.font.size = Pt(11)

    def save(self, output_path: str) -> bool:
        """
        儲存 PPT 到指定路徑

        Args:
            output_path: 輸出檔案路徑

        Returns:
            是否成功儲存
        """
        if not self.prs:
            return False

        try:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(output_path)
            return True
        except Exception:
            return False
