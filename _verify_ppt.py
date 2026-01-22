# -*- coding: utf-8 -*-
"""Verify PPT contains picture shape"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

ppt_path = Path('output/ppt') / '產業專業類' / '李文正_(韋能能源)_CV.pptx'
if ppt_path.exists():
    prs = Presentation(str(ppt_path))
    slide = prs.slides[0]
    print('=== Generated PPT Shapes ===')
    for i, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        name = shape.name
        left = shape.left / 914400  # EMU to inches
        top = shape.top / 914400
        width = shape.width / 914400
        height = shape.height / 914400
        type_str = f'{shape_type.name} ({int(shape_type)})'
        print(f'{i}: {name} - Type: {type_str}')
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f'   Photo: left={left:.2f}in, top={top:.2f}in, {width:.2f}x{height:.2f}in')

    pic_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    print(f'')
    print(f'Total PICTURE shapes: {pic_count}')
    if pic_count > 0:
        print('SUCCESS: Photo is present in the PPT!')
else:
    print('PPT file not found')
